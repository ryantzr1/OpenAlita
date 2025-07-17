import os
import shutil
from agent.gaia_agent import GAIAAgent
import pandas as pd
from PIL import Image
import wave
from pptx import Presentation
import zipfile
import json
from docx import Document
import numpy as np

TEST_DIR = "GAIA/2023"
DUMMY_FILES = [
        "example.txt",
        "example.csv",
        "example.xlsx",
        "example.pdb",
        "example1.pdf",
        "example2.pdf",
        "example.jpg",
        "example.png",
        "example.wav",
        "example.pptx",
        "example.zip",
        "example.jsonld",
        "example.py",
        "example.docx"
        "example.MOV"
    ]

def create_dummy_files():
    os.makedirs(TEST_DIR, exist_ok=True)

    # .txt
    with open(os.path.join(TEST_DIR, "example.txt"), "w") as f:
        f.write("This is a dummy text file for GAIA.")

    # .csv
    df = pd.DataFrame({"a": [1, 2], "b": [3, 4]})
    df.to_csv(os.path.join(TEST_DIR, "example.csv"), index=False)

    # .xlsx
    df.to_excel(os.path.join(TEST_DIR, "example.xlsx"), index=False)

    # .pdb
    with open(os.path.join(TEST_DIR, "example.pdb"), "w") as f:
        f.write("HEADER    DUMMY PDB FILE\nATOM      1  N   ALA A   1")

    # .jpg
    img = Image.new("RGB", (100, 100), color="blue")
    img.save(os.path.join(TEST_DIR, "example.jpg"))

    # .png
    img.save(os.path.join(TEST_DIR, "example.png"))
    
    # .pdf
    #1
    img.save(os.path.join(TEST_DIR, "example1.pdf"), "PDF")
    #2
    from reportlab.pdfgen import canvas
    pdf_path = os.path.join(TEST_DIR, "example2.pdf")
    c = canvas.Canvas(pdf_path)
    c.drawString(100, 750, "Dummy PDF content for testing.")
    c.save()

    # .wav
    wav_path = os.path.join(TEST_DIR, "example.wav")
    with wave.open(wav_path, 'w') as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(44100)
        wf.writeframes(b'\x00\x00' * 44100)

    # .pptx
    ppt = Presentation()
    slide_layout = ppt.slide_layouts[0]
    slide = ppt.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Dummy PPTX Slide"
    ppt.save(os.path.join(TEST_DIR, "example.pptx"))

    # .zip
    zip_path = os.path.join(TEST_DIR, "example.zip")
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        zipf.writestr("dummy.txt", "This is a file inside a zip archive.")

    # .jsonld
    jsonld_data = {
        "@context": "http://schema.org",
        "@type": "Person",
        "name": "Dummy Name",
        "jobTitle": "Software Agent"
    }
    with open(os.path.join(TEST_DIR, "example.jsonld"), "w") as f:
        json.dump(jsonld_data, f, indent=2)

    # .py
    with open(os.path.join(TEST_DIR, "example.py"), "w") as f:
        f.write("print('Hello from dummy Python script')\n")

    # .docx
    doc = Document()
    doc.add_paragraph("This is a dummy Word document for testing.")
    doc.save(os.path.join(TEST_DIR, "example.docx"))

    # .mov (creating a minimal MOV file)
    mov_path = os.path.join(TEST_DIR, "example.MOV")
    # Create a very small video file (just header basically)
    with open(mov_path, 'wb') as f:
        f.write(b'\x00\x00\x00\x20ftypqt  \x00\x00\x00\x01qt  ')
    

    print("✅ All dummy files created successfully!")

def test_DUMMY_file_loading():
    agent = GAIAAgent(gaia_files_dir=TEST_DIR)
    print(f"\n🔍 [TEST] Testing GAIA Agent File Loading...\n")

    for file_name in DUMMY_FILES:
        print(f"\n--- Testing file: {file_name} ---")
        result = agent._load_file_content(file_name)
        if result is None:
            print("❌ Failed to load or unsupported file format.")
        else:
            print("✅ File loaded successfully. Preview:")
            print(result[:500])

def clean_up():
    print("\n🧹 Cleaning up dummy files...")
    for file_name in DUMMY_FILES:
        file_path = os.path.join(TEST_DIR, file_name)
        if os.path.exists(file_path):
            os.remove(file_path)
            print(f"Removed: {file_path}")
    print("✅ Cleanup complete. Other files remain untouched.")

def test_GAIA_file_loading():
    validation_dir = os.path.join(TEST_DIR, "validation")
    test_dir = os.path.join(TEST_DIR, "test")

    agent_validation = GAIAAgent(gaia_files_dir=validation_dir)
    agent_test = GAIAAgent(gaia_files_dir=test_dir)

    print(f"\n🔍 [TEST] Testing all files in {validation_dir} ...\n")
    if os.path.exists(validation_dir):
        for fname in os.listdir(validation_dir):
            file_path = os.path.join(validation_dir, fname)
            if os.path.isfile(file_path) and fname.endswith(('.m4a')):
                print(f"\n--- Testing file: {fname} ---")
                result = agent_validation._load_file_content(fname)
                if result is None:
                    print("❌ Failed to load or unsupported file format.")
                else:
                    print("✅ File loaded successfully. Preview:")
                    print(result[:200])
    else:
        print(f"Directory not found: {validation_dir}")

    print(f"\n🔍 [TEST] Testing all files in {test_dir} ...\n")
    if os.path.exists(test_dir):
        for fname in os.listdir(test_dir):
            file_path = os.path.join(test_dir, fname)
            if os.path.isfile(file_path) and fname.endswith(('.m4a', '.mp3')):
                print(f"\n--- Testing file: {fname} ---")
                result = agent_test._load_file_content(fname)
                if result is None:
                    print("❌ Failed to load or unsupported file format.")
                else:
                    print("✅ File loaded successfully. Preview:")
                    print(result[:200])
    else:
        print(f"Directory not found: {test_dir}")
    

if __name__ == "__main__":
    create_dummy_files()
    test_DUMMY_file_loading()
    clean_up()
    test_GAIA_file_loading()
